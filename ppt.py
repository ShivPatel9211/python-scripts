from pptx import Presentation
from pptx.util import Inches
prs = Presentation()

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Adding a Bullet Slide'
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'
img_path='image1.jpg'
left = Inches(1)
top=Inches(2.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
prs.save('test.pptx')